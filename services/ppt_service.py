import pptx
import os

class PPTService:
    def __init__(self):
        self.supported_formats = ['.pptx', '.ppt']
        
    def extract_text(self, ppt_file):
        """提取PPT文本内容"""
        if not os.path.exists(ppt_file):
            return None
            
        try:
            prs = pptx.Presentation(ppt_file)
            text_content = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text_content.append("\n".join(slide_text))
                    
            return "\n\n".join(text_content)
        except Exception as e:
            print(f"PPT解析错误: {str(e)}")
            return None